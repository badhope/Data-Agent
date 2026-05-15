"""
PPT生成模块
基于python-pptx的演示文稿生成功能
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from typing import List, Dict, Optional
import io
from datetime import datetime


class PPTGenerator:
    """基于python-pptx的PPT生成器"""

    def __init__(self):
        self.presentation = None
        self.current_slide = None
        self._setup_chinese_support()

    def _setup_chinese_support(self):
        """设置中文支持"""
        from pptx.oxml.xmlchemy import OxmlElement
        
        def set_element_text(element, text):
            """设置元素文本，支持中文"""
            if text:
                element.text = text
        
        self._set_text = set_element_text

    def create_presentation(self, title: str = "演示文稿", author: str = "DataAgent"):
        """创建新演示文稿"""
        self.presentation = Presentation()
        # 使用英文避免编码问题
        self.presentation.core_properties.author = "DataAgent" if author else "DataAgent"
        self.presentation.core_properties.title = "Presentation" if title else "Presentation"
        self.presentation.core_properties.created = datetime.now()

    def add_title_slide(self, title: str, subtitle: str = ""):
        """添加标题页"""
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle

        self.current_slide = slide
        return slide

    def add_content_slide(self, title: str, content: List[str],
                          layout_index: int = 1):
        """添加内容页"""
        slide_layout = self.presentation.slide_layouts[layout_index]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()

            for i, point in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = "• " + point if not point.startswith(('•', '-', '1.', '2.')) else point
                p.level = 0

        self.current_slide = slide
        return slide

    def add_two_column_slide(self, title: str, left_content: List[str],
                            right_content: List[str]):
        """添加双栏内容页"""
        slide_layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
        else:
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5),
                                              Inches(8), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True

        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5),
                                          Inches(4.25), Inches(5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        for i, point in enumerate(left_content):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = "• " + point if not point.startswith(('•', '-')) else point

        right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5),
                                           Inches(4.25), Inches(5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        for i, point in enumerate(right_content):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = "• " + point if not point.startswith(('•', '-')) else point

        self.current_slide = slide
        return slide

    def add_image_slide(self, title: str, image_path: str = None,
                       description: str = ""):
        """添加图片页"""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        if image_path:
            try:
                slide.shapes.add_picture(
                    image_path,
                    Inches(1),
                    Inches(2),
                    width=Inches(8)
                )
            except Exception as e:
                pass

        if description:
            desc_box = slide.shapes.add_textbox(
                Inches(1),
                Inches(6),
                Inches(8),
                Inches(1)
            )
            desc_frame = desc_box.text_frame
            p = desc_frame.paragraphs[0]
            p.text = description
            p.alignment = PP_ALIGN.CENTER

        self.current_slide = slide
        return slide

    def add_table_slide(self, title: str, headers: List[str],
                       rows: List[List[str]]):
        """添加表格页"""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        rows_count = len(rows) + 1
        cols_count = len(headers)

        table = slide.shapes.add_table(
            rows_count,
            cols_count,
            Inches(0.5),
            Inches(1.8),
            Inches(9),
            Inches(0.5 * rows_count)
        ).table

        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            from pptx.dml.color import RGBColor
            cell.fill.fore_color.rgb = RGBColor(0, 51, 102)

        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)

        self.current_slide = slide
        return slide

    def add_section_slide(self, section_title: str):
        """添加章节分隔页"""
        slide_layout = self.presentation.slide_layouts[2]
        slide = self.presentation.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = section_title

        self.current_slide = slide
        return slide

    def add_blank_slide(self):
        """添加空白页"""
        slide_layout = self.presentation.slide_layouts[6]
        slide = self.presentation.slides.add_slide(slide_layout)
        self.current_slide = slide
        return slide

    def add_chart_slide(self, title: str, data: List[List], chart_type: str = "bar"):
        """
        添加图表页
        
        Args:
            title: 图表标题
            data: 图表数据，格式为 [[标签1, 值1], [标签2, 值2], ...]
            chart_type: 图表类型，支持 bar（柱状图）、line（折线图）、pie（饼图）
        
        Returns:
            创建的幻灯片对象
        """
        slide_layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title
        else:
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
        
        # 创建图表
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        
        chart_data = CategoryChartData()
        
        # 提取数据
        if len(data) > 0 and isinstance(data[0], list):
            headers = data[0]
            chart_data.categories = headers[1:] if len(headers) > 1 else [str(i) for i in range(len(data)-1)]
            
            # 添加数据系列
            for row in data[1:]:
                if len(row) > 0:
                    name = row[0]
                    values = row[1:]
                    chart_data.add_series(name, values)
        else:
            # 简单格式：[[标签, 值], ...]
            chart_data.categories = [str(row[0]) for row in data]
            chart_data.add_series(title, [row[1] for row in data])
        
        # 设置图表位置和大小
        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
        
        # 创建图表
        if chart_type.lower() == "line":
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
        elif chart_type.lower() == "pie":
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            ).chart
        else:
            # 默认柱状图
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
        
        # 设置图表标题
        chart.chart_title.text_frame.text = title
        
        self.current_slide = slide
        return slide

    def set_slide_background(self, color: tuple = (255, 255, 255)):
        """设置幻灯片背景颜色"""
        if self.current_slide:
            background = self.current_slide.background
            fill = background.fill
            fill.solid()
            from pptx.dml.color import RGBColor
            fill.fore_color.rgb = RGBColor(*color)

    def save(self, output_path: str):
        """保存演示文稿"""
        if self.presentation:
            import io
            if isinstance(output_path, io.BytesIO):
                self.presentation.save(output_path)
            else:
                self.presentation.save(output_path)
            return True
        return False

    def save_to_bytes(self) -> bytes:
        """保存演示文稿为字节流"""
        import io
        output = io.BytesIO()
        if self.presentation:
            try:
                self.presentation.save(output)
                return output.getvalue()
            except Exception as e:
                import logging
                logging.error(f"PPT save error: {e}")
                return b''
        return b''

    def save_to_buffer(self) -> bytes:
        """保存演示文稿到内存缓冲区（别名方法）"""
        return self.save_to_bytes()

    def save_to_file(self, file_path: str) -> bool:
        """
        保存演示文稿到文件
        
        Args:
            file_path: 输出文件路径
            
        Returns:
            是否保存成功
        """
        try:
            # 确保目录存在
            import os
            dir_path = os.path.dirname(file_path)
            if dir_path:
                os.makedirs(dir_path, exist_ok=True)
            
            # 保存文件
            with open(file_path, 'wb') as f:
                f.write(self.save_to_bytes())
            
            return True
        except Exception as e:
            import logging
            logging.error(f"Save to file error: {e}")
            return False

    def get_bytes(self) -> bytes:
        """获取二进制内容"""
        if not self.presentation:
            return b""

        import io
        buffer = io.BytesIO()
        try:
            # 确保所有文本都是字符串类型
            self._sanitize_text()
            self.presentation.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        except Exception as e:
            import logging
            logging.error(f"PPT get_bytes error: {e}")
            return b""

    def _sanitize_text(self):
        """清理所有幻灯片中的文本，确保编码正确"""
        if not self.presentation:
            return
        
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                # 确保文本是字符串
                                run.text = str(run.text)
                elif hasattr(shape, "text"):
                    if shape.text:
                        shape.text = str(shape.text)

    def get_slide_count(self) -> int:
        """获取幻灯片数量"""
        return len(self.presentation.slides) if self.presentation else 0


class PPTTemplate:
    """PPT模板管理器"""

    TEMPLATES = {
        "business": {
            "name": "商业报告",
            "description": "专业的商业分析报告模板",
            "slides": ["封面", "目录", "概述", "数据分析", "结论建议"],
            "color_scheme": {
                "primary": (0, 51, 102),
                "secondary": (0, 102, 153),
                "accent": (255, 153, 0)
            }
        },
        "academic": {
            "name": "学术报告",
            "description": "学术研究演示模板",
            "slides": ["标题", "研究背景", "方法", "结果", "讨论", "结论"],
            "color_scheme": {
                "primary": (51, 51, 51),
                "secondary": (102, 102, 102),
                "accent": (0, 102, 204)
            }
        },
        "meeting": {
            "name": "会议纪要",
            "description": "简洁的会议演示模板",
            "slides": ["会议主题", "议程", "讨论要点", "行动计划", "下次会议"],
            "color_scheme": {
                "primary": (51, 51, 51),
                "secondary": (102, 102, 102),
                "accent": (204, 0, 0)
            }
        },
        "proposal": {
            "name": "项目提案",
            "description": "项目方案展示模板",
            "slides": ["项目概述", "目标", "方案", "时间表", "预算", "风险评估"],
            "color_scheme": {
                "primary": (0, 102, 102),
                "secondary": (0, 153, 153),
                "accent": (255, 153, 0)
            }
        },
        "weekly_report": {
            "name": "周报",
            "description": "工作周报模板",
            "slides": ["本周总结", "本周完成", "本周进展", "下周计划", "问题与建议"],
            "color_scheme": {
                "primary": (0, 102, 51),
                "secondary": (0, 153, 76),
                "accent": (255, 153, 0)
            }
        }
    }

    @classmethod
    def get_template(cls, template_id: str) -> Dict:
        """获取模板配置"""
        return cls.TEMPLATES.get(template_id, cls.TEMPLATES["business"])

    @classmethod
    def list_templates(cls) -> List[Dict]:
        """列出所有模板"""
        return [
            {"id": key, **value}
            for key, value in cls.TEMPLATES.items()
        ]

    @classmethod
    def generate_from_template(
        cls,
        template_id: str,
        title: str,
        content: Dict[str, List[str]]
    ) -> bytes:
        """使用模板生成PPT"""
        template = cls.get_template(template_id)
        generator = PPTGenerator()
        generator.create_presentation(title)

        generator.add_title_slide(title, template["name"])

        for section_title in template["slides"][1:]:
            section_key = section_title.replace(" ", "_").lower()
            section_content = content.get(section_key, [])

            if section_content:
                generator.add_content_slide(section_title, section_content)

        return generator.get_bytes()

    @classmethod
    def generate(cls, template_id: str, data: Dict[str, any]) -> Dict[str, any]:
        """
        使用模板生成PPT并保存到文件
        
        Args:
            template_id: 模板ID (business, academic, meeting, proposal, weekly_report)
            data: 包含PPT内容的数据字典
            
        Returns:
            包含生成结果的字典
        """
        try:
            template = cls.get_template(template_id)
            generator = PPTGenerator()
            
            # 创建演示文稿
            title = data.get("title", "演示文稿")
            generator.create_presentation(title)
            
            # 添加标题页
            subtitle = data.get("subtitle", "")
            generator.add_title_slide(title, subtitle)
            
            # 根据模板类型添加内容页
            if template_id == "business":
                # 商业报告模板
                if "overview" in data:
                    generator.add_content_slide("概述", data["overview"])
                if "data" in data:
                    generator.add_chart_slide("数据分析", data["data"], chart_type="bar")
                if "conclusion" in data:
                    generator.add_content_slide("结论建议", data["conclusion"])
                    
            elif template_id == "academic":
                # 学术报告模板
                if "background" in data:
                    generator.add_content_slide("研究背景", data["background"])
                if "method" in data:
                    generator.add_content_slide("研究方法", data["method"])
                if "results" in data:
                    generator.add_content_slide("研究结果", data["results"])
                if "discussion" in data:
                    generator.add_content_slide("讨论", data["discussion"])
                if "conclusion" in data:
                    generator.add_content_slide("结论", data["conclusion"])
                    
            elif template_id == "meeting":
                # 会议纪要模板
                if "agenda" in data:
                    generator.add_content_slide("议程", data["agenda"])
                if "discussion" in data:
                    generator.add_content_slide("讨论要点", data["discussion"])
                if "action" in data:
                    generator.add_content_slide("行动计划", data["action"])
                    
            elif template_id == "proposal":
                # 项目提案模板
                if "overview" in data:
                    generator.add_content_slide("项目概述", data["overview"])
                if "goals" in data:
                    generator.add_content_slide("目标", data["goals"])
                if "plan" in data:
                    generator.add_content_slide("方案", data["plan"])
                if "timeline" in data:
                    generator.add_content_slide("时间表", data["timeline"])
                if "budget" in data:
                    generator.add_content_slide("预算", data["budget"])
                    
            elif template_id == "weekly_report":
                # 周报模板
                if "summary" in data:
                    generator.add_content_slide("本周总结", data["summary"])
                if "completed" in data:
                    generator.add_content_slide("本周完成", data["completed"])
                if "progress" in data:
                    generator.add_content_slide("本周进展", data["progress"])
                if "next_week" in data:
                    generator.add_content_slide("下周计划", data["next_week"])
                if "issues" in data:
                    generator.add_content_slide("问题与建议", data["issues"])
            
            # 通用内容添加
            if "contents" in data:
                for item in data["contents"]:
                    if isinstance(item, dict) and "title" in item and "content" in item:
                        generator.add_content_slide(item["title"], item["content"])
            
            # 生成输出路径
            import os
            output_dir = "output"
            os.makedirs(output_dir, exist_ok=True)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(output_dir, f"{template_id}_report_{timestamp}.pptx")
            
            # 保存文件
            generator.save_to_file(output_path)
            
            return {
                "success": True,
                "message": "PPT生成成功",
                "output_path": output_path,
                "template": template_id,
                "title": title
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }
