"""
文档处理服务
提供文档摘要、格式化、PPT生成、PDF解析等功能
"""
import io
import json
import re
from typing import Dict, List, Optional, Any
from datetime import datetime
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from services.llm_service import LLMService

llm_service = LLMService()


class DocumentService:
    """文档处理服务类"""

    def __init__(self):
        self.llm = llm_service

    async def summarize_document(
        self,
        text: str,
        method: str = "extractive",
        max_length: int = 200,
        num_sentences: int = 5
    ) -> Dict[str, Any]:
        """生成文档摘要

        Args:
            text: 文档内容
            method: 摘要方法(extractive/abstractive)
            max_length: 最大长度
            num_sentences: 句子数量

        Returns:
            包含摘要结果的字典
        """
        if not text or not text.strip():
            return {"summary": "请提供文档内容"}

        try:
            if method == "extractive":
                summary = self._extractive_summary(text, num_sentences, max_length)
            else:
                summary = await self._abstractive_summary(text, max_length)

            return {"summary": summary, "method": method}
        except Exception as e:
            return {"summary": f"摘要生成失败: {str(e)}", "error": True}

    def _extractive_summary(self, text: str, num_sentences: int, max_length: int) -> str:
        """抽取式摘要 - 提取最重要的句子"""
        sentences = re.split(r'[。！？\n]+', text)
        sentences = [s.strip() for s in sentences if s.strip()]

        if len(sentences) <= num_sentences:
            return text[:max_length] + "..." if len(text) > max_length else text

        scored_sentences = []
        for i, sentence in enumerate(sentences):
            score = 0
            if len(sentence) > 10:
                score += 1
            if i < 3:
                score += 2
            if i > len(sentences) - 3:
                score += 1
            words = set(sentence)
            score += min(len(words) / 20, 5)
            scored_sentences.append((sentence, score))

        scored_sentences.sort(key=lambda x: x[1], reverse=True)
        selected = scored_sentences[:num_sentences]
        selected.sort(key=lambda x: sentences.index(x[0]))

        summary = "".join([s[0] for s in selected])
        return summary[:max_length] + "..." if len(summary) > max_length else summary

    async def _abstractive_summary(self, text: str, max_length: int) -> str:
        """生成式摘要 - 使用AI生成摘要"""
        prompt = f"""请对以下文档生成简洁的摘要，控制在{max_length}字以内：

{text[:3000]}

摘要："""

        try:
            response = await self.llm.generate(prompt)
            return response.strip()[:max_length]
        except Exception:
            return self._extractive_summary(text, 3, max_length)

    async def summarize_structured(
        self,
        text: str,
        document_type: str = "general"
    ) -> Dict[str, Any]:
        """生成结构化摘要"""
        prompt = f"""请对以下文档进行结构化摘要分析，输出JSON格式：

文档类型：{document_type}
文档内容：
{text[:3000]}

请输出包含以下字段的JSON：
{{
    "title": "文档标题",
    "summary": "简短摘要(100字内)",
    "key_points": ["要点1", "要点2", "要点3"],
    "keywords": ["关键词1", "关键词2"],
    "category": "文档类别"
}}

JSON输出："""

        try:
            response = await self.llm.generate(prompt)
            json_match = re.search(r'\{[\s\S]*\}', response)
            if json_match:
                result = json.loads(json_match.group())
                return result
            return {"summary": response, "raw": True}
        except Exception as e:
            return {"summary": f"结构化摘要生成失败: {str(e)}", "error": True}

    async def generate_meeting_minutes(
        self,
        text: str,
        meeting_date: Optional[str] = None,
        output_format: str = "markdown"
    ) -> str:
        """生成会议纪要"""
        date_str = meeting_date or datetime.now().strftime("%Y-%m-%d")

        prompt = f"""请根据以下会议记录生成专业的会议纪要：

会议日期：{date_str}
会议记录：
{text[:3000]}

请生成包含以下部分的会议纪要（使用Markdown格式）：
1. 会议基本信息
2. 会议主题
3. 讨论内容
4. 决议事项
5. 待办事项（用- [ ]格式）
6. 下次会议安排

会议纪要："""

        try:
            minutes = await self.llm.generate(prompt)
            if output_format == "html":
                minutes = minutes.replace("\n", "<br>")
            return minutes
        except Exception as e:
            return f"会议纪要生成失败: {str(e)}"

    async def generate_meeting_ppt(
        self,
        meeting_text: str,
        meeting_date: Optional[str] = None
    ) -> bytes:
        """生成会议纪要PPT"""
        minutes = await self.generate_meeting_minutes(meeting_text, meeting_date)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
        title_frame = title.text_frame
        p = title_frame.paragraphs[0]
        p.text = "会议纪要"
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        date_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(0.8))
        date_frame = date_box.text_frame
        p = date_frame.paragraphs[0]
        p.text = meeting_date or datetime.now().strftime("%Y年%m月%d日")
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        sections = self._parse_sections(minutes)
        for section_title, section_content in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            header = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            header_frame = header.text_frame
            p = header_frame.paragraphs[0]
            p.text = section_title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 64, 175)

            content = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
            content_frame = content.text_frame
            content_frame.word_wrap = True
            p = content_frame.paragraphs[0]
            p.text = section_content[:500] if len(section_content) > 500 else section_content
            p.font.size = Pt(18)

        pptx_stream = io.BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream.read()

    def _parse_sections(self, text: str) -> List[tuple]:
        """解析文本中的章节"""
        sections = []
        lines = text.split("\n")
        current_title = ""
        current_content = []

        for line in lines:
            if re.match(r'^#+\s+', line) or re.match(r'^\d+\.\s+\w+', line):
                if current_title:
                    sections.append((current_title, "\n".join(current_content)))
                current_title = re.sub(r'^#+\s+', '', line).strip()
                current_content = []
            else:
                current_content.append(line)

        if current_title:
            sections.append((current_title, "\n".join(current_content)))

        if not sections:
            sections = [("会议纪要", text)]

        return sections

    async def generate_report(
        self,
        work_items: List[Dict],
        report_type: str = "weekly",
        author: str = "DataAgent"
    ) -> str:
        """生成工作报告"""
        items_text = "\n".join([f"- {item.get('content', item)}" for item in work_items])

        prompt = f"""请根据以下工作内容生成{report_type}报告：

作者：{author}
报告类型：{report_type}
工作内容：
{items_text}

请生成包含以下部分的{report_type}报告：
1. 报告概述
2. 本期完成工作
3. 工作成果
4. 遇到的问题与解决方案
5. 下期工作计划

报告："""

        try:
            report = await self.llm.generate(prompt)
            return report
        except Exception as e:
            return f"报告生成失败: {str(e)}"

    async def generate_report_ppt(
        self,
        work_items: List[Dict],
        report_type: str = "weekly"
    ) -> bytes:
        """生成工作报告PPT"""
        report = await self.generate_report(work_items, report_type)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        title = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
        title_frame = title.text_frame
        p = title_frame.paragraphs[0]
        p.text = f"{report_type.capitalize()} Report"
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        sections = self._parse_sections(report)
        for i, (section_title, section_content) in enumerate(sections[:8]):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            header = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            header_frame = header.text_frame
            p = header_frame.paragraphs[0]
            p.text = section_title
            p.font.size = Pt(28)
            p.font.bold = True

            content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
            content_frame = content.text_frame
            content_frame.word_wrap = True
            p = content_frame.paragraphs[0]
            p.text = section_content[:800] if len(section_content) > 800 else section_content
            p.font.size = Pt(16)

        pptx_stream = io.BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream.read()

    async def extract_todos(self, text: str) -> Dict[str, Any]:
        """提取待办事项"""
        prompt = f"""请从以下文本中提取待办事项，输出JSON格式：

{text[:2000]}

请输出JSON格式：
{{
    "todos": [
        {{"content": "待办事项1", "deadline": "截止日期或null", "priority": "high/medium/low"}},
        ...
    ],
    "meetings": [
        {{"content": "会议内容", "date": "日期", "attendees": ["人员"]}},
        ...
    ]
}}

JSON输出："""

        try:
            response = await self.llm.generate(prompt)
            json_match = re.search(r'\{[\s\S]*\}', response)
            if json_match:
                result = json.loads(json_match.group())
                return result
            return {"todos": [], "raw": response}
        except Exception as e:
            return {"todos": [], "error": str(e)}

    async def format_text(
        self,
        text: str,
        format_level: str = "standard"
    ) -> Dict[str, Any]:
        """格式化文本"""
        prompt = f"""请对以下文本进行格式化处理：

原始文本：
{text}

格式化要求：{format_level}
- 中英文之间添加空格
- 标点符号标准化
- 保持原文格式

格式化结果："""

        try:
            formatted = await self.llm.generate(prompt)
            return {"formatted": formatted, "original_length": len(text)}
        except Exception as e:
            return {"formatted": text, "error": str(e)}

    def get_ppt_templates(self) -> List[Dict[str, str]]:
        """获取PPT模板列表"""
        return [
            {"id": "business", "name": "商务风格", "description": "适合正式商务汇报"},
            {"id": "simple", "name": "简约风格", "description": "简洁大方，突出内容"},
            {"id": "creative", "name": "创意风格", "description": "适合创意展示"},
            {"id": "academic", "name": "学术风格", "description": "适合学术报告"}
        ]

    async def generate_ppt_from_content(
        self,
        title: str,
        content: Dict[str, List[str]],
        template: str = "business"
    ) -> bytes:
        """根据内容生成PPT"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(2))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        for section_title, items in content.items():
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            header = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
            header_frame = header.text_frame
            p = header_frame.paragraphs[0]
            p.text = section_title
            p.font.size = Pt(36)
            p.font.bold = True

            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.8))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, item in enumerate(items[:8]):
                if i == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(24)
                p.space_after = Pt(12)

        pptx_stream = io.BytesIO()
        prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream.read()

    async def manage_citations(
        self,
        text: str = "",
        action: str = "format",
        citation_id: str = "",
        style: str = "gbt"
    ) -> Dict[str, Any]:
        """管理引用格式"""
        if action == "format":
            prompt = f"""请将以下参考文献转换为{style}格式：

参考文献：
{text}

请输出转换后的引用格式："""

            try:
                result = await self.llm.generate(prompt)
                return {"result": result, "style": style}
            except Exception as e:
                return {"result": "", "error": str(e)}
        elif action == "list":
            styles = {
                "gbt": "GB/T 7714",
                "apa": "APA",
                "mla": "MLA",
                "chicago": "Chicago"
            }
            return {
                "result": f"当前支持的引用格式：{', '.join(styles.values())}",
                "available_styles": list(styles.keys())
            }
        return {"result": "", "action": action}

    async def parse_pdf(self, pdf_bytes: bytes, extract_tables: bool = False) -> Dict[str, Any]:
        """解析PDF文档"""
        try:
            pdf_file = io.BytesIO(pdf_bytes)

            with pdfplumber.open(pdf_file) as pdf:
                text_content = []
                tables_content = []

                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"[第{page_num + 1}页]\n{page_text}")

                    if extract_tables:
                        tables = page.extract_tables()
                        for table in tables:
                            if table:
                                tables_content.append({
                                    "page": page_num + 1,
                                    "data": table
                                })

                result = {
                    "page_count": len(pdf.pages),
                    "text": "\n\n".join(text_content),
                    "has_tables": len(tables_content) > 0,
                    "tables": tables_content if extract_tables else None
                }

                if not text_content[0] if text_content else False:
                    return {"error": "无法解析PDF内容"}

                return result

        except Exception as e:
            return {"error": f"PDF解析失败: {str(e)}"}

    async def search_pdf(self, content: Dict, query: str) -> List[Dict]:
        """在PDF内容中搜索"""
        text = content.get("text", "")
        if not query or not text:
            return []

        query_words = query.lower().split()
        text_lower = text.lower()

        results = []
        for word in query_words:
            if word in text_lower:
                index = text_lower.find(word)
                start = max(0, index - 50)
                end = min(len(text), index + len(word) + 50)
                snippet = text[start:end]
                results.append({
                    "query": word,
                    "snippet": "..." + snippet + "...",
                    "position": index
                })

        return results

    async def generate_outline(
        self,
        topic: str,
        document_type: str = "general",
        depth: int = 3
    ) -> Dict[str, Any]:
        """生成文章提纲"""
        prompt = f"""请为以下主题生成{document_type}类型的文章提纲：

主题：{topic}
文档类型：{document_type}
大纲深度：{depth}级

请输出JSON格式：
{{
    "title": "文章标题",
    "outline": [
        {{"level": 1, "content": "一级标题", "sub": [...]}},
        ...
    ],
    "estimated_length": "预估字数"
}}

JSON输出："""

        try:
            response = await self.llm.generate(prompt)
            json_match = re.search(r'\{[\s\S]*\}', response)
            if json_match:
                result = json.loads(json_match.group())
                return result
            return {"outline": response, "title": topic}
        except Exception as e:
            return {"error": str(e)}

    async def generate_outline_markdown(
        self,
        topic: str,
        document_type: str = "general",
        depth: int = 3
    ) -> str:
        """生成Markdown格式提纲"""
        outline_data = await self.generate_outline(topic, document_type, depth)

        if "error" in outline_data:
            return f"提纲生成失败: {outline_data['error']}"

        markdown = f"# {outline_data.get('title', topic)}\n\n"

        def add_items(items, level=1):
            result = []
            for item in items:
                if isinstance(item, dict):
                    prefix = "#" + "#" * level
                    result.append(f"{prefix} {item.get('content', '')}")
                    if 'sub' in item:
                        result.extend(add_items(item['sub'], level + 1))
                elif isinstance(item, str):
                    prefix = "#" + "#" * level
                    result.append(f"{prefix} {item}")
            return result

        if 'outline' in outline_data:
            for item in add_items(outline_data['outline']):
                markdown += item + "\n"

        return markdown

    def get_outline_templates(self) -> List[Dict[str, str]]:
        """获取提纲模板列表"""
        return [
            {"id": "general", "name": "通用文档", "description": "适用于一般文章写作"},
            {"id": "report", "name": "工作报告", "description": "适合工作汇报、述职"},
            {"id": "proposal", "name": "方案策划", "description": "适合项目方案、营销策划"},
            {"id": "academic", "name": "学术论文", "description": "适合学术研究论文"},
            {"id": "research", "name": "调研报告", "description": "适合市场调研、行业分析"}
        ]
