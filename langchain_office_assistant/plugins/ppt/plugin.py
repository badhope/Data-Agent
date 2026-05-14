from typing import List, Dict, Any, Optional
from langchain_core.tools import tool
from langchain_office_assistant.plugins.base import BasePlugin
from langchain_office_assistant.utils.logger import get_logger
from pptx import Presentation
from pptx.util import Inches, Pt

logger = get_logger(__name__)

class PPTPlugin(BasePlugin):
    name = "ppt"
    description = "PPT生成插件 - 模板支持、图表嵌入"

    def __init__(self):
        super().__init__()

    def initialize(self, config: Dict) -> None:
        self.config = config
        logger.info(f"PPTPlugin initialized")

    def get_tools(self) -> List:
        return [create_ppt, add_slide, add_chart_to_ppt, save_ppt]

    async def execute(self, tool_name: str, **kwargs) -> Any:
        tools_map = {
            "create_ppt": create_ppt,
            "add_slide": add_slide,
            "add_chart_to_ppt": add_chart_to_ppt,
            "save_ppt": save_ppt,
        }

        if tool_name not in tools_map:
            return f"❌ Tool not found: {tool_name}"

        tool_func = tools_map[tool_name]
        return tool_func.invoke(kwargs)

_ppt_instances = {}

@tool
def create_ppt(title: str, template: Optional[str] = None) -> str:
    """Create a new PowerPoint presentation."""
    try:
        if template:
            prs = Presentation(template)
        else:
            prs = Presentation()

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        instance_id = f"ppt_{id(prs)}"
        _ppt_instances[instance_id] = prs

        return f"✅ PowerPoint presentation created!\n\n📌 Title: {title}\n🆔 Instance ID: {instance_id}"
    except Exception as e:
        return f"❌ Failed to create PPT: {str(e)}"

@tool
def add_slide(
    instance_id: str,
    slide_type: str = "title",
    title: Optional[str] = None,
    content: Optional[str] = None,
    bullet_points: Optional[List[str]] = None
) -> str:
    """Add a new slide to an existing presentation."""
    try:
        if instance_id not in _ppt_instances:
            return f"❌ Presentation not found: {instance_id}"

        prs = _ppt_instances[instance_id]

        layout_map = {
            "title": 0,
            "content": 1,
            "section": 2,
            "title_only": 5,
            "blank": 6,
        }

        layout_idx = layout_map.get(slide_type, 1)
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        if title:
            if slide.shapes.title:
                slide.shapes.title.text = title

        if content:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if not shape.text or shape == slide.shapes.title:
                        continue
                    shape.text = content

        if bullet_points:
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    tf = shape.text_frame
                    tf.clear()
                    for point in bullet_points:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0

        return f"✅ Slide added successfully!\n\n📄 Slide Type: {slide_type}\n📌 Title: {title or 'N/A'}"
    except Exception as e:
        return f"❌ Failed to add slide: {str(e)}"

@tool
def add_chart_to_ppt(
    instance_id: str,
    chart_type: str = "bar",
    title: str = "Chart",
    categories: List[str] = None,
    values: List[float] = None
) -> str:
    """Add a chart to an existing presentation slide."""
    try:
        if instance_id not in _ppt_instances:
            return f"❌ Presentation not found: {instance_id}"

        prs = _ppt_instances[instance_id]
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        categories = categories or ["A", "B", "C", "D"]
        values = values or [10, 20, 30, 40]

        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series("Values", values)

        x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(5)

        chart_types = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE,
            "line": XL_CHART_TYPE.LINE,
            "area": XL_CHART_TYPE.AREA,
        }

        chart_type_enum = chart_types.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        slide.shapes.add_chart(chart_type_enum, x, y, cx, cy, chart_data)

        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = title

        return f"✅ Chart added successfully!\n\n📊 Chart Type: {chart_type}\n📌 Title: {title}"
    except Exception as e:
        return f"❌ Failed to add chart: {str(e)}"

@tool
def save_ppt(instance_id: str, file_path: str) -> str:
    """Save a PowerPoint presentation to a file."""
    try:
        if instance_id not in _ppt_instances:
            return f"❌ Presentation not found: {instance_id}"

        prs = _ppt_instances[instance_id]
        prs.save(file_path)

        del _ppt_instances[instance_id]

        return f"✅ Presentation saved!\n\n📄 File: {file_path}"
    except Exception as e:
        return f"❌ Failed to save PPT: {str(e)}"
